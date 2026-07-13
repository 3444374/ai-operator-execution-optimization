from pathlib import Path

from pptx import Presentation

src = Path("projects/opening_defense_20260712/exports/opening_defense_20260712_201257.pptx")
dst = Path("projects/opening_defense_20260712/exports/opening_defense_20260712_fixed.pptx")

prs = Presentation(src)
w, h = prs.slide_width, prs.slide_height

# Slide 7: the chapter-title textbox cloned from the template is wider than
# the slide. Constrain it to the safe content width.
slide = prs.slides[6]
shape = slide.shapes[1]
shape.width = w - shape.left * 2

# Slide 16: the references body is slightly taller than the slide. Keep the
# same top position but constrain its height to the safe lower margin.
slide = prs.slides[15]
shape = slide.shapes[2]
shape.height = h - shape.top - 450000

prs.save(dst)
print(dst)
