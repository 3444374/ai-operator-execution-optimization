from pptx import Presentation

path = "projects/opening_defense_20260712/exports/opening_defense_20260712_200753.pptx"
prs = Presentation(path)
w, h = prs.slide_width, prs.slide_height

for idx in [6, 15]:
    slide = prs.slides[idx]
    print("slide", idx + 1, "size", int(w), int(h))
    for j, shape in enumerate(slide.shapes):
        text = shape.text if hasattr(shape, "text") else ""
        if not text.strip():
            continue
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        out = left < 0 or top < 0 or left + width > w or top + height > h
        print(j, out, int(left), int(top), int(width), int(height), text.replace("\n", " ")[:100])
