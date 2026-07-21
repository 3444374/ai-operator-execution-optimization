"""
开题答辩PPT生成脚本 v3
- 严格遵循学校模板布局：编号块 #1588C7 + 标题 28pt + #E2E2E2 分隔线
- 汇报人/导师留空占位
- 研究内容页图文混排，文字精简扼要
- 20页结构：封面·目录 + 4段内容 + 致谢
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# === Paths ===
TEMPLATE = "opening/templates/模板.pptx"
OUTPUT = "opening/slides/opening_defense_20260720_v3.pptx"
FIG = "figures"
FIG_CROP = "figures/ppt_cropped"  # Cropped data charts (no titles/captions)

def img_path(name, cropped=False):
    """Resolve figure path. Cropped versions for data charts."""
    base = FIG_CROP if cropped else FIG
    return os.path.join(base, name)

# === Canvas (16:9) ===
SW = Inches(13.333)
SH = Inches(7.5)

# === Template Colors (from 模板.pptx analysis) ===
C_BLOCK_BG  = RGBColor(0x15, 0x88, 0xC7)  # Section number block fill
C_BLOCK_FG  = RGBColor(0xFF, 0xFF, 0xFF)  # Section number text (white)
C_TITLE     = RGBColor(0x00, 0x00, 0x00)  # Title text (black / TEXT_1)
C_BODY      = RGBColor(0x00, 0x00, 0x00)  # Body text (black)
C_MUTED     = RGBColor(0x6B, 0x72, 0x80)  # Muted/secondary
C_SEP       = RGBColor(0xE2, 0xE2, 0xE2)  # Separator line
C_BULLET    = RGBColor(0x24, 0x4C, 0x89)  # Bullet marker squares
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# === Four-section accent colors (段落标识色) ===
SEC_COLORS = {
    "01": RGBColor(0x2F, 0x6F, 0xEB),  # Blue
    "02": RGBColor(0xF9, 0x73, 0x16),  # Orange
    "03": RGBColor(0x16, 0xA3, 0x4A),  # Green
    "04": RGBColor(0x7C, 0x3A, 0xED),  # Purple
}

# === Template Header Positions (from 模板.pptx measurement) ===
HDR_BLOCK_L = 0.60   # Section number block left
HDR_BLOCK_T = 1.07   # Section number block top
HDR_BLOCK_W = 0.80   # Section number block width
HDR_BLOCK_H = 0.71   # Section number block height
HDR_TITLE_L = 1.65   # Section title left
HDR_TITLE_T = 1.12   # Section title top
HDR_SEP_Y   = 1.83   # Separator line top
HDR_SEP_W   = 11.90  # Separator line width
CONTENT_T   = 2.10   # Content area starts below

# =====================================================
# HELPERS
# =====================================================

def delete_all_slides(prs):
    """Remove all existing slides from the presentation."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def add_rect(slide, l, t, w, h, fill_color, text="", fs=14, bold=False, fg=C_WHITE, align=PP_ALIGN.CENTER, font_name='Microsoft YaHei'):
    """Add a filled rectangle with optional centered text."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    rect.line.fill.background()
    if text:
        tf = rect.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(fs)
        p.font.bold = bold
        p.font.color.rgb = fg
        p.alignment = align
        p.font.name = font_name
    return rect


def add_textbox(slide, l, t, w, h, text, fs=14, bold=False, color=C_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name='Microsoft YaHei'):
    """Add a text box."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    tb.text_frame.auto_size = None
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.font.name = font_name
    return tb


def add_bullets(slide, l, t, w, h, items, fs=14, spacing=6, font_name='Microsoft YaHei'):
    """Add bulleted list. Each bullet prefixed with a #244C89 square marker (template style)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Template-style bullet: small dark-navy square + text
        p.text = f"■ {item}"  # ■
        p.font.size = Pt(fs)
        p.font.color.rgb = C_BULLET if i == 0 else C_BODY
        p.font.name = font_name
        p.space_after = Pt(spacing)
    # Set color of the ■ symbol explicitly
    for p in tf.paragraphs:
        run = p.runs[0] if p.runs else None
        if run:
            run.font.color.rgb = C_BULLET
    return tb


def add_img_fit(slide, path, l, t, max_w, max_h):
    """Add image, fitting within max_w × max_h box while preserving aspect ratio."""
    if not os.path.exists(path):
        return None
    pic = slide.shapes.add_picture(path, Inches(l), Inches(t))
    w_in = pic.width / 914400
    h_in = pic.height / 914400
    scale = min(max_w / w_in, max_h / h_in, 1.0)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    # Center within the box
    if pic.width / 914400 < max_w:
        pic.left = Inches(l + (max_w - pic.width / 914400) / 2)
    return pic


def set_slide_number(slide, num):
    """Set slide number in SLIDE_NUMBER placeholder."""
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 13:  # SLIDE_NUMBER
            ph.text = str(num)


def add_notes(slide, speech, defense):
    """Add dual-section speaker notes: 汇报讲稿 + 答辩备注."""
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.clear()

    def ap(text, bold=False, fs=11):
        p = tf.add_paragraph()
        p.text = text
        p.font.bold = bold
        p.font.size = Pt(fs)

    ap("汇报讲稿：", bold=True)
    if speech:
        ap(speech, fs=10)
    ap("")
    ap("答辩备注：", bold=True)
    if defense:
        ap(defense, fs=10)


# =====================================================
# STANDARD SECTION HEADER (matching template)
# =====================================================

def add_section_header(slide, section_num, title):
    """Template-consistent section header:
       - Blue block (#1588C7) with section number at (0.60, 1.07)
       - Title text at (1.65, 1.12), 28pt bold
       - Gray separator line (#E2E2E2) at y=1.83
    """
    # Number block
    add_rect(slide, HDR_BLOCK_L, HDR_BLOCK_T, HDR_BLOCK_W, HDR_BLOCK_H,
             C_BLOCK_BG, section_num, fs=32, bold=True, fg=C_BLOCK_FG)
    # Title
    add_textbox(slide, HDR_TITLE_L, HDR_TITLE_T, 10.0, 0.55,
                title, fs=28, bold=True, color=C_TITLE)
    # Separator
    add_rect(slide, HDR_BLOCK_L, HDR_SEP_Y, HDR_SEP_W, 0.02, C_SEP)


# =====================================================
# LAYOUT PATTERNS
# =====================================================

def layout_cover(prs):
    """Cover slide: Layout 0 (title slide), leave author/advisor blank."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Title (placeholder 0)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            r = ph.text_frame.paragraphs[0].add_run()
            r.text = "数据库 AI 负载的\n执行优化与调度研究"
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.color.rgb = C_TITLE
        elif ph.placeholder_format.idx == 1:
            ph.text = ""
            r = ph.text_frame.paragraphs[0].add_run()
            r.text = "Daft/Arrow数据组织 · Ray执行调度 · GPU推理服务 · 持久化协同"
            r.font.size = Pt(14)
            r.font.color.rgb = C_MUTED
            r.alignment = PP_ALIGN.CENTER
        elif ph.placeholder_format.type == 16:  # DATE
            ph.text = "2026年7月"

    # Presenter info — leave blank as requested
    add_textbox(slide, 1.71, 4.47, 9.0, 0.8,
                "报告人：____________    指导老师：____________",
                fs=18, color=C_TITLE, align=PP_ALIGN.LEFT)

    return slide


def layout_toc(prs, items):
    """TOC slide: Layout 1, template-style simple text list."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Title placeholder → "目录"
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            r = ph.text_frame.paragraphs[0].add_run()
            r.text = "目录"
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.color.rgb = C_TITLE
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item
                p.font.size = Pt(20)
                p.font.bold = (i == 0)
                p.font.color.rgb = C_TITLE
                p.space_after = Pt(12)

    return slide


def layout_section_divider(prs, num, title):
    """Section divider: Layout 2 (blank), large section number + title + accent line."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    color = SEC_COLORS.get(num, C_BLOCK_BG)

    # Large section number
    add_textbox(slide, 0.55, 1.2, 2.0, 1.3, num, fs=72, bold=True, color=color)
    # Title
    add_textbox(slide, 2.5, 1.5, 9.0, 0.8, title, fs=36, bold=True, color=C_TITLE)
    # Section accent line
    add_rect(slide, 2.5, 2.5, 8.0, 0.04, color)
    # Thin secondary
    add_rect(slide, 2.5, 2.7, 4.0, 0.01, color)

    return slide


def layout_content_img(prs, section, title, bullets, img_name, img_on_left=False, conclusion="", cropped=False):
    """Content + image: template header, narrow text left + large image right."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    color = SEC_COLORS.get(section, C_BLOCK_BG)
    add_section_header(slide, section, title)

    text_l = 0.60
    text_w = 4.80    # narrower text column
    img_l = 5.80     # image starts earlier
    img_w = 7.00     # wider image area
    img_h = 3.80     # max image height
    img_t = CONTENT_T + 0.05

    if img_on_left:
        text_l, img_l = img_l, text_l
        if not img_on_left:
            img_l = 5.80

    # Bullets — larger font
    add_bullets(slide, text_l, CONTENT_T, text_w, img_h, bullets, fs=15)

    # Image — fit within box, preserve AR
    add_img_fit(slide, img_path(img_name, cropped), img_l, img_t, img_w, img_h)

    if conclusion:
        add_rect(slide, 0.60, 6.90, 0.08, 0.45, color)
        add_textbox(slide, 0.85, 6.90, 11.0, 0.45, conclusion, fs=12, bold=True, color=color)

    return slide


def layout_img_top(prs, section, title, bullets, img_name, conclusion="", cropped=False):
    """Image-dominant layout: compact bullets top, large image below (constrained to slide)."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    color = SEC_COLORS.get(section, C_BLOCK_BG)
    add_section_header(slide, section, title)

    # Compact bullets at top
    add_bullets(slide, 0.60, CONTENT_T, 12.0, 0.85, bullets, fs=13, spacing=4)

    # Image area: from 3.0in to 6.2in bottom (3.2in available), full width with margins
    img_t = CONTENT_T + 1.0
    img_max_h = 6.20 - img_t - 0.05
    add_img_fit(slide, img_path(img_name, cropped), 0.55, img_t, 12.20, img_max_h)

    if conclusion:
        add_rect(slide, 0.60, 6.90, 0.08, 0.45, color)
        add_textbox(slide, 0.85, 6.90, 11.0, 0.45, conclusion, fs=12, bold=True, color=color)

    return slide


def layout_text_only(prs, section, title, bullets, conclusion=""):
    """Text-only content slide with template header. Used sparingly."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    color = SEC_COLORS.get(section, C_BLOCK_BG)
    add_section_header(slide, section, title)

    add_bullets(slide, 0.60, CONTENT_T, 12.0, 3.50, bullets, fs=15, spacing=8)

    if conclusion:
        add_rect(slide, 0.60, 6.90, 0.08, 0.45, color)
        add_textbox(slide, 0.85, 6.90, 11.0, 0.45, conclusion, fs=12, bold=True, color=color)

    return slide


def layout_two_col(prs, section, title, left_title, left_items, right_title, right_items, conclusion=""):
    """Two-column layout for comparison (feasibility + innovation, etc.)."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    color = SEC_COLORS.get(section, C_BLOCK_BG)
    add_section_header(slide, section, title)

    # Left column
    add_textbox(slide, 0.60, CONTENT_T, 5.80, 0.35, left_title, fs=16, bold=True, color=color)
    add_bullets(slide, 0.60, CONTENT_T + 0.40, 5.80, 3.30, left_items, fs=14, spacing=5)

    # Right column
    add_textbox(slide, 6.80, CONTENT_T, 5.80, 0.35, right_title, fs=16, bold=True, color=color)
    add_bullets(slide, 6.80, CONTENT_T + 0.40, 5.80, 3.30, right_items, fs=14, spacing=5)

    if conclusion:
        add_rect(slide, 0.60, 6.90, 0.08, 0.45, color)
        add_textbox(slide, 0.85, 6.90, 11.0, 0.45, conclusion, fs=11, bold=True, color=color)

    return slide


def layout_timeline(prs, section, title, phases, conclusion=""):
    """Horizontal timeline flow. Uses blank layout with section header."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    color = SEC_COLORS.get(section, C_BLOCK_BG)
    add_section_header(slide, section, title)

    n = len(phases)
    box_w = 11.5 / n - 0.12
    row_y = CONTENT_T + 0.20

    for i, (label, desc) in enumerate(phases):
        x = 0.55 + i * (box_w + 0.12)
        # Phase box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(row_y), Inches(box_w), Inches(1.15))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER
        # Arrow between boxes
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           Inches(x + box_w + 0.01), Inches(row_y + 0.45),
                                           Inches(0.10), Inches(0.22))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()
        # Description below
        add_textbox(slide, x, row_y + 1.30, box_w, 2.30, desc, fs=12, color=C_BODY)

    if conclusion:
        add_rect(slide, 0.60, 6.90, 0.08, 0.45, color)
        add_textbox(slide, 0.85, 6.90, 11.0, 0.45, conclusion, fs=11, bold=True, color=color)

    return slide


def layout_end(prs):
    """Ending slide: Layout 3, '谢谢各位老师'."""
    slide = prs.slides.add_slide(prs.slide_layouts[3])

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            r = ph.text_frame.paragraphs[0].add_run()
            r.text = "谢谢各位老师"
            r.font.size = Pt(36)
            r.font.bold = True
            r.font.color.rgb = C_TITLE
            ph.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        elif ph.placeholder_format.idx == 1:
            ph.text = ""
            r = ph.text_frame.paragraphs[0].add_run()
            r.text = "恳请批评指正"
            r.font.size = Pt(20)
            r.font.color.rgb = C_MUTED
            ph.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


# =====================================================
# BUILD DECK
# =====================================================

def build(prs):
    n = 0  # slide counter

    # ====== COVER ======
    n += 1
    s = layout_cover(prs)
    add_notes(s,
        "各位老师好，我汇报的题目是'数据库 AI 负载的执行优化与调度研究'。本课题关注数据库成为AI工作负载入口后，数据从数据库表出发经分布式数据处理、GPU推理再写回存储这条端到端链路中的上游调度优化问题。",
        "如被问题目是否太大：研究对象是固定执行路径 PostgreSQL→Daft/Arrow→Ray→vLLM→写回，在可控链路上做上游调度策略消融研究。不修改vLLM内部、不改造Ray核心、不做GPU kernel优化。")
    set_slide_number(s, n)

    # ====== TOC ======
    n += 1
    s = layout_toc(prs, [
        "01  研究背景与问题",
        "02  预研基础与动机证据",
        "03  研究目标与内容",
        "04  可行性与计划",
    ])
    add_notes(s,
        "汇报分为四个部分：研究背景与核心问题、预研实验与动机证据、研究目标与两项内容设计、可行性分析与进度安排。",
        "四段结构遵循学校模板，每段有过渡页和对应内容页。")
    set_slide_number(s, n)

    # ========================================
    # SECTION 01: 研究背景与问题
    # ========================================

    n += 1
    s = layout_section_divider(prs, "01", "研究背景与问题")
    add_notes(s, "首先进入第一部分：研究背景与问题。", "")
    set_slide_number(s, n)

    n += 1
    s = layout_text_only(prs, "01", "数据库成为AI Workload入口",
        ["Snowflake Cortex AISQL（SIGMOD 2026）：六类AI SQL算子，约40%查询消耗58%+总执行时间——AI算子已成SQL一等公民",
         "BigQuery ML/AI、Oracle AI Vector Search：ML.GENERATE_TEXT、VECTOR_EMBEDDING等函数，支持Gemini/Claude/Llama",
         "PostgreSQL生态：pgvector向量存储 + pgai外部vectorizer + PostgresML，形成可观测的开源替代方案"],
        conclusion="AI SQL算子已是工业现实——数据库正从管理结构化数据扩展为AI数据处理入口。")
    add_notes(s,
        "云数据库厂商已在SQL中提供AI算子。Snowflake Cortex AISQL论文发表于SIGMOD 2026。BigQuery和Oracle同样在SQL中嵌入AI函数。PostgreSQL生态提供了可观测的开源替代方案。",
        "Snowflake等系统作为工业背景引用，不声称它们使用Ray或不复现其内部实现——它们证明需求真实性，不能直接作为消融实验baseline。")
    set_slide_number(s, n)

    n += 1
    s = layout_content_img(prs, "01", "AI算子执行链路与传统SQL的差异",
        ["传统SQL算子：在数据库执行器内部完成（scan→filter→join→aggregate），数据不离开数据库进程",
         "AI算子链路更长：表数据→Daft/Arrow batch组织→Ray调度执行→GPU推理服务→结果汇聚→写回",
         "端到端性能不只由GPU推理速度决定——数据组织方式、提交节奏、模型服务队列状态、结果写回批量共同影响"],
        img_name="architecture/system_architecture_ai_data_execution.png",
        img_on_left=True,
        conclusion="AI算子的外部执行链路需要独立优化——上游数据如何组织为请求、以什么节奏发送、如何根据模型服务状态调节并发。")
    add_notes(s,
        "传统SQL算子在数据库内部完成。AI算子不同——数据经过外部数据处理层、调度层、GPU推理，最后写回。端到端性能受上游多个阶段影响。",
        "研究对象不是模型kernel优化，而是数据从数据库出发到进入推理引擎之前的上游调度链路。Snowflake证实重要性但作为闭源系统不可拆分——这正是开题动机。")
    set_slide_number(s, n)

    n += 1
    s = layout_content_img(prs, "01", "研究现状与核心研究问题",
        ["DB4AI路线（GaussML/Smart/NeurDB）：将ML嵌入数据库内核，优化止于数据库进程边界",
         "GPU推理服务（vLLM/Orca/Sarathi-Serve）：优化GPU内存与批处理调度，将数据来源抽象为「输入请求」",
         "AI数据存储（TurboVecDB/Delta Lake/Lance）：优化存储格式与写入路径，不研究上游数据调度"],
        img_name="architecture/research_gap_three_islands.png",
        conclusion="核心问题：数据从数据库表到GPU推理再到写回的上游调度——三个成熟方向之间的连接处——尚未被作为系统性优化目标研究。")
    add_notes(s,
        "现有工作分三个方向。DB4AI止于数据库进程边界。vLLM等优化GPU侧但不关心数据从哪来。AI存储优化写入但不研究上游调度。三个方向各有大量CCF-A论文，但连接处是空白。本课题聚焦数据组织与执行调度两项策略。",
        "研究缺口表述：不声称'没人做过AI SQL'，而是'缺少面向本课题执行路径的可控阶段画像与优化验证'。")
    set_slide_number(s, n)

    # ========================================
    # SECTION 02: 预研基础与动机证据
    # ========================================

    n += 1
    s = layout_section_divider(prs, "02", "预研基础与动机证据")
    add_notes(s, "第二部分展示已完成的预研实验和动机证据。", "")
    set_slide_number(s, n)

    n += 1
    s = layout_img_top(prs, "02", "GPU-backed AI_EMBED端到端链路",
        ["实验环境：本地PG18.4同构预演 + RTX 5070 12GB + sentence-transformers/all-MiniLM-L6-v2 + CUDA endpoint",
         "完整链路：PostgreSQL → Arrow RecordBatch → Ray/Python → CUDA embedding endpoint → writeback，阶段计时均可独立观测"],
        img_name="02_gpu_stage_latency_stack.png", cropped=True,
        conclusion="真实GPU-backed AI_EMBED链路已跑通——阶段计时方法可行，端到端链路可观测、可拆分、可消融。")
    add_notes(s,
        "AI_EMBED预研已完成。完整链路包含DB读取、Arrow batch、Ray执行、GPU推理和写回。各阶段计时独立可拆分。",
        "主动标注：当前是PG18.4本地预演，不是PG18.3内部平台结果。AI_EMBED仅作为预研验证，论文主体实验在vLLM + AI_COMPLETE平台。")
    set_slide_number(s, n)

    n += 1
    s = layout_content_img(prs, "02", "预研关键结果：三组动机信号",
        ["① Batch粒度：1024行fine vs coalesced→端到端差异37.5×——调用粒度是一阶执行成本，非GPU kernel优化问题",
         "② Writeback成本：4096行JSON 1.567s，pgvector(384) 0.897s——写回已是可见成本，当前不主导端到端",
         "③ Ray价值边界：单endpoint下Ray actor(3.62s)不优于Python(3.42s)；双endpoint降至2.86s——调度价值依赖并行条件"],
        img_name="10_e2e_operator_writeback_breakdown.png", cropped=True,
        conclusion="三组信号共同指向：batch粒度和调度策略是上游优先调优方向；写回已量化作为端到端保护约束。")
    add_notes(s,
        "三组关键结果。37.5倍差异来自调用粒度而非GPU优化。写回已量化。Ray不天然优于Python——多endpoint才体现路由价值。",
        "边界：(1) 37.5×来自调用粒度，非GPU kernel优化；(2) 双endpoint在同一GPU上不代表多GPU结论；(3) 写回不是主瓶颈但已量化备查。")
    set_slide_number(s, n)

    n += 1
    s = layout_content_img(prs, "02", "vLLM + AI_COMPLETE Baseline",
        ["发现①：固定行数≠固定计算量——batch=8时token跨度13.9×，batch=128时token P95=26,678，行数是计算量的弱代理",
         "发现②：改为按token量攒批（token-budget），约束P95至~6,141，同时吞吐接近固定行方案——按计算量组织数据可行",
         "发现③：共享vLLM下无界提交伤害并发——前台作业端到端从4.9s恶化到11.4s（2.3×），证明需要in-flight控制",
         "Daft管线已验证：开销<0.1s，Rust内核+Arrow零拷贝——后续多模态实验可复用同一套pipeline代码"],
        img_name="b24_local_vllm_interference_sweep_small_job.png", cropped=True,
        conclusion="三项核心发现构成论文可写证据。Daft接入为多模态泛化验证铺路——同一套策略代码，只需替换数据列类型。")
    add_notes(s,
        "7月18-19日完成vLLM baseline和Daft接入。三项发现构成论文正文可写证据。图中展示shared-vLLM下不同K_max对前台作业的冲击——无界提交导致2.3×恶化。Daft管线开销已验证<0.1s，且Daft支持@daft.cls GPU UDF，后续多模态实验只需替换数据列类型（prompt→image），复用同一套策略代码。",
        "澄清：vLLM baseline是本地rehearsal。Queue-adaptive已实现但不如静态K_max——这是开题后需攻克的技术难点。开题阶段只展示动机成立，不声称已完成。")
    set_slide_number(s, n)

    # ========================================
    # SECTION 03: 研究目标与内容
    # ========================================

    n += 1
    s = layout_section_divider(prs, "03", "研究目标与内容")
    add_notes(s, "第三部分介绍研究目标和两项策略设计。", "")
    set_slide_number(s, n)

    n += 1
    s = layout_img_top(prs, "03", "研究目标与总体框架",
        ["总体目标：以AI_COMPLETE（生成式LLM推理）为主场景，构建Daft/Ray端到端实验链路，优化上游数据组织与提交控制两项策略",
         "主链路：PostgreSQL → Daft（数据引擎，Rust核心+Arrow零拷贝）→ Ray actor（异构池+去中心化协调）→ vLLM（continuous batching，不修改内部）→ 写回",
         "两项策略 + 多模态泛化验证 + 算子代价估计（补充）：研究内容一（数据组织）、研究内容二（提交控制）、耦合验证（独立拼接 vs 联合grid search）"],
        img_name="architecture/cross_layer_method_framework.png",
        conclusion="核心思路：数据如何组织、以什么节奏发送、如何响应模型服务状态——三项决策共同影响端到端效率。")
    add_notes(s,
        "总体目标：基于Daft/Ray端到端实验链路，优化数据组织策略和提交控制策略。vLLM为部署平台不修改，Ray为架构设计空间。",
        "不能把题目解释成'优化Ray'。真正问题是数据库AI workload的数据组织、服务感知调度和持久化协同。")
    set_slide_number(s, n)

    n += 1
    s = layout_content_img(prs, "03", "研究内容一：按计算量组织数据，而非按行数",
        ['问题：固定"每批N行"简单，但LLM每行token数差异巨大（实测13.9×跨度）——长文本拖着短文本一起等，短文本被白白拖慢',
         '做法① Token-budget：不数行数，数token数——每攒够一定token量就发一批。短行多的批次行数自然多，长行多的批次行数自然少',
         '做法② Length-align：按token长度排序后分组，让每批内的行计算量相近——减少"一人拖全组"的straggler效应',
         '做法③ Prefix-aware：共享system prompt的行合并提交，让vLLM复用KV-cache，避免重复计算同一段前缀'],
        img_name="architecture/data_organization_strategy_mechanism.png",
        conclusion="核心思路：用token量替代行数作为攒批单位——让计算量决定批次边界，而非机械地数行。引擎参数（Daft into_batches/repartition）+策略决策共同优化。")
    add_notes(s,
        "研究内容一的核心问题是：固定行数不等于固定计算量。LLM每行token数差异巨大（实测13.9×跨度），所以按token量而不是按行数来组织batch是更合理的做法。三种策略从三个角度优化：token-budget控制每批总计算量，length-align减少批内straggler，prefix-aware让vLLM复用KV-cache。通过Ray actor异构化实现——不同策略对应不同actor类型。",
        '引擎级参数（Daft into_batches、batch_size等）与策略级决策共同构成优化空间——不是简单的"选一个batch size"。')
    set_slide_number(s, n)  # p13

    n += 1
    s = layout_content_img(prs, "03", "研究内容二：让请求发送节奏响应模型服务状态",
        ["问题：固定每隔N行或每隔T秒发一次请求，不管GPU忙不忙——GPU空闲时浪费算力，GPU繁忙时雪上加霜",
         "做法① Queue-adaptive flush：每个Ray actor独立看vLLM队列（当前在跑几个、在排几个），队列空了就多发，队列堵了就少发——用实时状态替代固定间隔",
         '做法② K_max动态调整：不预设"最多同时发K个"，而是由queue-adaptive行为自然形成上限——随观测窗口自动升降',
         "做法③ Actor pool分池：按请求特征分池——短token走快池、长token走慢池、共享prefix的走亲和池，各行按特征路由到不同提交策略"],
        img_name="architecture/runtime_strategy_rule_table.png",
        conclusion="利用Ray actor的状态记忆+异步能力，每个actor独立做决策——去中心化，不引入全局协调瓶颈。")
    add_notes(s,
        "研究内容二关注提交控制——什么时候把攒好的请求发给vLLM。固定间隔的问题是：GPU空闲时你还在等，GPU忙时你又雪上加霜。Queue-adaptive让每个actor独立看vLLM队列状态决定发不发：队列空就发，队列堵就等。K_max不再预设固定值，由adaptive行为自然形成。Actor pool按请求特征分池路由。",
        "当前queue-adaptive已实现但不如静态K_max=8——开题后需改进。若3轮改进仍不超越则降级。但K_max必要性已通过shared-vLLM干扰实验证明（2.3×恶化）。")
    set_slide_number(s, n)  # p14

    n += 1
    s = layout_text_only(prs, "03", "耦合验证与实验设计",
        ["验证问题：两项策略分别调优后拼接，是否等于联合调优？还是需要联合搜索才能找到最优组合？",
         "方法：独立搜索最优batching配置 → 独立搜索最优submission配置 → 拼接独立最优 → 联合grid search（batching × submission）→ 比较差距",
         "判定：联合显著优于拼接 → 需要联合调优（交互效应强）；两者接近 → 分层独立优化即可（策略正交）",
         "无论哪种结果，都不改变课题的核心贡献——上游数据组织+提交控制两项策略的设计空间表征与实验验证。"],
        conclusion="写回使用PostgreSQL+pgvector（COPY+deferred index）作为工程baseline，不作为独立方法贡献——仅在端到端评价中检查写回是否吞噬上游收益。")
    add_notes(s,
        "耦合验证是AGENTS.md写死的核心实验。方法：分别独立搜索RC1最优 → 独立搜索RC2最优 → 拼接 → 与联合grid search对比。如果联合显著优于拼接，说明两项策略有交互效应需要联合调优。如果接近，说明策略正交可以分层独立优化。无论哪种结果都不改变核心贡献。写回作为保护约束纳入端到端评价。",
        "耦合验证必须回答——不能只分别报告两项策略各自最优，而不说明它们之间是否需要联合调优。")
    set_slide_number(s, n)  # p15

    n += 1
    s = layout_timeline(prs, "03", "实验阶段路线图", [
        ("前置阶段", "vLLM baseline\nDaft接入\n已完成"),
        ("第一阶段 8月", "数据组织消融\nToken-budget\nLength-align\nPrefix-aware"),
        ("第二阶段 9月", "提交控制消融\nQueue-adaptive\nK_max动态\nPool路由"),
        ("第三阶段 9-10月", "耦合验证\n独立拼接 vs\n联合grid search"),
        ("补充阶段 10月", "多模态泛化验证\n写回 guardrail\n算子代价估计"),
    ], conclusion="前置条件已完成，后续四阶段每阶段有明确验证目标和反证条件。")
    add_notes(s,
        "实验分四阶段。前置已全部完成。后续按数据组织→提交控制→耦合验证顺序推进。多模态泛化验证为正文实验(§5.3)，写回作为 guardrail，算子代价估计为补充讨论。每阶段有明确对照和反证。",
        "前置完成证明方案可执行。如被问'4个月能否完成'：前置已满足，每阶段实验规模可控，脚本已模板化。")
    set_slide_number(s, n)

    # ========================================
    # SECTION 04: 可行性与计划
    # ========================================

    n += 1
    s = layout_section_divider(prs, "04", "可行性与计划")
    add_notes(s, "第四部分分析可行性和进度安排。", "")
    set_slide_number(s, n)

    n += 1
    s = layout_two_col(prs, "04", "可行性与创新点",
        "可行性依据",
        ["真实GPU-backed链路已跑通——AI_EMBED预研 + AI_COMPLETE baseline，阶段计时方法成熟",
         "Daft数据引擎接入已验证——管线开销<0.1s，与手动Arrow路径同数量级",
         "风险降级路径明确：自适应3轮不超越则降级为K_max必要性论证；文本消融完成前不启动多模态"],
        "预期创新点",
        ["创新点①：AI workload感知的上游动态数据组织策略——在Ray侧设计token-budget/length-align/prefix-aware三种batching policy",
         "创新点②：Ray actor去中心化自适应提交策略——每个actor独立观测vLLM队列状态自主决定flush时机",
         "创新点③：AI数据流持久化瓶颈判定——通过耦合验证与sink对比为全链路优化提供写入侧边界条件"],
        conclusion="方法已通过预研验证；硬件边界不影响核心方法；风险有明确触发条件和降级路径。创新点为预期，需后续实验验证。")
    add_notes(s,
        "可行性基于链路已跑通、Daft已验证、vLLM baseline已建立。硬件单机单GPU不影响核心方法验证。预期创新点需要后续实验验证。",
        "主动说明硬件边界：单机单GPU。核心方法（数据组织、in-flight控制、routing）不依赖分布式，单机可验证。")
    set_slide_number(s, n)

    n += 1
    s = layout_text_only(prs, "04", "进度安排",
        ["2026.07  开题报告 ✅ | vLLM baseline建立 ✅ | Daft接入 ✅ | token-tail revision ✅",
         "2026.08  研究内容一消融：Token-budget / Length-align / Prefix-aware vs 静态baseline",
         "2026.09  研究内容二消融：Queue-adaptive / K_max动态 / Pool路由 vs 固定K_max + 耦合验证启动",
         "2026.10  多模态泛化验证 + 写回 guardrail + 算子代价估计（补充讨论）",
         "2026.11  整理统一方法，补齐消融与反证实验，形成可复现脚本和完整CSV",
         "2026.12+ 论文实验、图表、正文撰写、答辩材料准备"],
        conclusion="按基线→消融→耦合验证顺序推进，多模态泛化验证+写回 guardrail 收尾，每阶段保留可反证条件。")
    add_notes(s,
        "进度从7月到12月。7月前置已完成。8-9月是核心实验窗口。10月多模态泛化验证。最后两个月整理写论文。每阶段有可反证条件。",
        "关键路径是8-9月消融实验——论文主体证据来源。如学校进度节点不同可按学院模板调整。")
    set_slide_number(s, n)

    # ====== SUMMARY ======
    n += 1
    s = layout_text_only(prs, "", "总结",
        ["场景真实：数据库已是AI workload入口（Snowflake/BigQuery/Oracle），AI算子外部执行链路需要独立优化",
         "证据初步成立：GPU-backed预研显示batch粒度37.5×差异 + writeback成本量化 + vLLM baseline三项核心发现",
         "后续聚焦：数据组织策略（token-budget/length-align/prefix-aware）+ 提交控制策略（queue-adaptive flush/pool路由）→ 耦合实验 → 多模态泛化验证",
         "边界清楚：不修改vLLM内部 · 不改造Ray核心 · 不做GPU kernel优化 · 不把CPU/fake当最终结论"],
        conclusion="题目有真实场景支撑，初步证据链成立，后续路线可执行且边界明确——恳请各位老师批评指正。")
    add_notes(s,
        "三点总结。场景真实——工业趋势。证据初步成立——关键信号。后续聚焦——两项策略加耦合验证。",
        "回答质疑优先回到阶段画像和不能声称的边界。最难的技术点是queue-adaptive控制器效果不如静态K_max——正在攻克。")
    set_slide_number(s, n)

    # ====== END ======
    n += 1
    s = layout_end(prs)
    add_notes(s, "我的汇报到此结束，恳请各位老师批评指正。谢谢！", "")
    set_slide_number(s, n)

    return n


# =====================================================
# MAIN
# =====================================================

if __name__ == "__main__":
    print(f"Opening template: {TEMPLATE}")
    prs = Presentation(TEMPLATE)
    print(f"Deleting {len(prs.slides)} existing slides...")
    delete_all_slides(prs)

    n = build(prs)

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Done! Generated {n} slides → {OUTPUT}")
